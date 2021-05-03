from pptx import Presentation
from hanspell import spell_checker
from hanspell.constants import CheckResult
import glob
import re
import sys


inputPath = "오탈자 확인\*.pptx"
sys.stdout = open('result.txt','w')

lis = glob.glob(inputPath)
print(*lis)


for path in glob.glob(inputPath):
    print(" ")
    print(path)
    f = open(path, "rb")
    prs = Presentation(f)
    # SPELLINGError = [] # 철자 오류
    # AMBIGUOUSError = [] #표준어 맞는지 모름
    # STATISTICAL_CORRECTIONError = [] #통계적으로 맞는 단어

    i = 1

    for slide in prs.slides:
        
        print("----------[#",i,"]------------")
        
        text = []

        for shape in slide.shapes:
            
            if not shape.has_text_frame:
                continue
            for run in shape.text_frame.paragraphs:
                run.text = re.sub(r'[^ ㄱ-ㅣ가-힣A-Za-z]', '', run.text)
                result = spell_checker.check(run.text)

                for key, value in result.words.items():
                    if value == 1:
                        print("WRONG_SPELLING")
                        print(result.original)
                        print(key)
                        
                    if value == 3:
                        print("AMBIGUOUS")
                        print(result.original)
                        print(key)
                        
                    if value == 4:
                        print("STATISTICAL_CORRECTION")
                        print(result.original)
                        print(key)
                        
        i += 1
        
        
        





            
            


            







